#!/usr/bin/env python3
"""
Generate the Technical Design Deck (.pptx) for this repository.

Outputs:
  - docs/diagrams/*.png  (rendered from docs/diagrams/*.dot via Graphviz 'dot')
  - docs/CounterUAS_RadarTracker_TDD.pptx
"""

from __future__ import annotations

import json
import os
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
DOCS_DIR = REPO_ROOT / "docs"
DIAGRAMS_DIR = DOCS_DIR / "diagrams"
OUT_PPTX = DOCS_DIR / "CounterUAS_RadarTracker_TDD.pptx"


BRAND_DARK = RGBColor(15, 23, 42)     # slate-900
BRAND_MUTED = RGBColor(71, 85, 105)   # slate-600
BRAND_ACCENT = RGBColor(2, 132, 199)  # sky-600


def _run(cmd: Sequence[str], cwd: Optional[Path] = None) -> None:
    subprocess.run(list(cmd), cwd=str(cwd) if cwd else None, check=True)


def render_diagrams() -> List[Path]:
    DIAGRAMS_DIR.mkdir(parents=True, exist_ok=True)
    rendered: List[Path] = []
    for dot in sorted(DIAGRAMS_DIR.glob("*.dot")):
        png = dot.with_suffix(".png")
        _run(["dot", "-Tpng", str(dot), "-o", str(png)], cwd=REPO_ROOT)
        rendered.append(png)
    return rendered


def load_tracker_config() -> dict:
    cfg_path = REPO_ROOT / "config" / "tracker_config.json"
    with cfg_path.open("r", encoding="utf-8") as f:
        return json.load(f)


@dataclass(frozen=True)
class SlideTheme:
    title_size_pt: int = 34
    body_size_pt: int = 18
    small_size_pt: int = 14
    margin_in: float = 0.6


def _set_run_font(run, size_pt: int, color: RGBColor, bold: bool = False) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_title(slide, title: str, theme: SlideTheme) -> None:
    left = Inches(theme.margin_in)
    top = Inches(0.3)
    width = Inches(13.333 - 2 * theme.margin_in)
    height = Inches(0.9)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    _set_run_font(r, theme.title_size_pt, BRAND_DARK, bold=True)


def add_subtitle(slide, text: str, theme: SlideTheme, *, top_in: float = 1.25) -> None:
    left = Inches(theme.margin_in)
    top = Inches(top_in)
    width = Inches(13.333 - 2 * theme.margin_in)
    height = Inches(0.8)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    _set_run_font(r, theme.body_size_pt, BRAND_MUTED, bold=False)


def add_bullets(slide, title: str, bullets: Sequence[str], theme: SlideTheme) -> None:
    add_title(slide, title, theme)
    left = Inches(theme.margin_in)
    top = Inches(1.35)
    width = Inches(13.333 - 2 * theme.margin_in)
    height = Inches(6.8 - 1.35)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(6)
        p.font.size = Pt(theme.body_size_pt)
        p.font.color.rgb = BRAND_DARK
        p.font.name = "Calibri"


def add_diagram_slide(
    slide,
    title: str,
    image_path: Path,
    theme: SlideTheme,
    caption: Optional[str] = None,
) -> None:
    add_title(slide, title, theme)

    # Reserve space for title and (optional) caption.
    left = Inches(theme.margin_in)
    top = Inches(1.2)
    width = Inches(13.333 - 2 * theme.margin_in)
    height = Inches(6.6 - (0.5 if caption else 0.0))

    pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    pic.line.color.rgb = RGBColor(226, 232, 240)  # slate-200

    if caption:
        c_left = Inches(theme.margin_in)
        c_top = Inches(7.1)
        c_width = Inches(13.333 - 2 * theme.margin_in)
        c_height = Inches(0.3)
        c = slide.shapes.add_textbox(c_left, c_top, c_width, c_height)
        tf = c.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = caption
        _set_run_font(r, theme.small_size_pt, BRAND_MUTED, bold=False)


def build_deck() -> None:
    theme = SlideTheme()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    cfg = load_tracker_config()

    # Slide 1: Title
    s = prs.slides.add_slide(blank)
    add_title(s, "Counter-UAS Radar Tracker", theme)
    add_subtitle(
        s,
        "Technical Design Documentation (TDD) — architecture, UML, sequences, and design principles",
        theme,
        top_in=1.35,
    )
    add_subtitle(
        s,
        f"Repo: {REPO_ROOT.name}  •  Output: {OUT_PPTX.name}",
        SlideTheme(body_size_pt=14, title_size_pt=34, small_size_pt=12, margin_in=0.6),
        top_in=2.05,
    )

    # Slide 2: What this system does
    s = prs.slides.add_slide(blank)
    add_bullets(
        s,
        "System overview",
        [
            "Purpose: ingest radar detections (per dwell) and output stabilized 3D tracks for downstream display / C2.",
            "Inputs: UDP SPDetectionMessage (range/azimuth/elevation + SNR/RCS/strength).",
            "Outputs: UDP TrackTableMessage (TrackUpdateMessage list) + optional binary logs.",
            "Primary concerns: real-time throughput, association correctness under clutter, track lifecycle management.",
        ],
        theme,
    )

    # Slide 3: Deployment / context
    s = prs.slides.add_slide(blank)
    add_diagram_slide(
        s,
        "Deployment / system context",
        DIAGRAMS_DIR / "deployment.png",
        theme,
        caption="UDP protocol and message structures are defined in `idl/messages.idl`.",
    )

    # Slide 4: Module / component architecture
    s = prs.slides.add_slide(blank)
    add_diagram_slide(
        s,
        "Component architecture (CMake targets)",
        DIAGRAMS_DIR / "component_architecture.png",
        theme,
        caption="Modules are linked as static libraries and composed by `cuas_pipeline` / `cuas_tracker`.",
    )

    # Slide 5: Threading model + runtime sequence
    s = prs.slides.add_slide(blank)
    add_diagram_slide(
        s,
        "Runtime sequence (threads and message handoff)",
        DIAGRAMS_DIR / "dwell_sequence.png",
        theme,
        caption="Receiver runs in its own thread; processing runs in a separate thread and is paced by `cyclePeriodMs`.",
    )

    # Slide 6: Data flow
    s = prs.slides.add_slide(blank)
    add_diagram_slide(
        s,
        "Per-dwell data flow (processing pipeline)",
        DIAGRAMS_DIR / "dwell_dataflow.png",
        theme,
    )

    # Slide 7: Data model
    s = prs.slides.add_slide(blank)
    add_bullets(
        s,
        "Core data model (wire + internal)",
        [
            "Detection: spherical measurement (range/az/el) plus features (strength/noise/SNR/RCS/microDoppler).",
            "Cluster: strength-weighted centroid of detections, converted to Cartesian (x,y,z) for filtering.",
            "Track state: 9D Cartesian IMM state [x,vx,ax,y,vy,ay,z,vz,az] with merged covariance and mode probabilities.",
            "TrackUpdateMessage: status + classification + spherical + Cartesian + quality/hits/misses/age.",
        ],
        theme,
    )

    # Slide 8: UML / class structure
    s = prs.slides.add_slide(blank)
    add_diagram_slide(
        s,
        "UML (high-level) — core classes and strategy interfaces",
        DIAGRAMS_DIR / "uml_class.png",
        theme,
        caption="Clustering and association are runtime-selectable via configuration (Strategy pattern).",
    )

    # Slide 9: Preprocessing & clustering
    s = prs.slides.add_slide(blank)
    add_bullets(
        s,
        "Preprocessing & clustering design",
        [
            "Preprocessor: hard-gate out-of-bounds detections (range, azimuth, elevation, SNR, RCS, strength).",
            "Clustering: groups detections into centroided targets to reduce association complexity.",
            "DBSCAN: distance in normalized (range,az,el) space using epsilons; noise points become singleton clusters.",
            "ClusterEngine assigns monotonically increasing clusterId and computes Cartesian centroid.",
        ],
        theme,
    )

    # Slide 10: Prediction (IMM)
    pred = cfg.get("prediction", {})
    imm = pred.get("imm", {})
    s = prs.slides.add_slide(blank)
    add_bullets(
        s,
        "Prediction — IMM filter",
        [
            "IMM models: CV, CA1, CA2, CTR1, CTR2 (all 9D Cartesian).",
            "Interaction (mixing) step uses transition matrix and current mode probabilities.",
            "Each model predicts its own (x,P); estimates are merged by mode probabilities.",
            f"Config: numModels={imm.get('numModels', 'n/a')}  •  initialModeProbabilities={imm.get('initialModeProbabilities', 'n/a')}",
        ],
        theme,
    )

    # Slide 11: Association
    assoc = cfg.get("association", {})
    s = prs.slides.add_slide(blank)
    add_bullets(
        s,
        "Association — measurement gating + assignment",
        [
            "Measurement space: Cartesian position (x,y,z) with H selecting position from the 9D state.",
            "Gating: Mahalanobis distance computed with innovation covariance S = HPHᵀ + R.",
            "Methods:",
            f"  - Mahalanobis NN: greedy nearest-neighbor with distanceThreshold={assoc.get('mahalanobis', {}).get('distanceThreshold', 'n/a')}",
            f"  - GNN: global assignment (reduced-cost greedy/Hungarian-like) with costThreshold={assoc.get('gnn', {}).get('costThreshold', 'n/a')}",
            f"  - JPDA: probabilistic association using clutterDensity={assoc.get('jpda', {}).get('clutterDensity', 'n/a')} and Pd={assoc.get('jpda', {}).get('detectionProbability', 'n/a')}",
        ],
        theme,
    )

    # Slide 12: Track lifecycle
    s = prs.slides.add_slide(blank)
    add_diagram_slide(
        s,
        "Track lifecycle (status transitions)",
        DIAGRAMS_DIR / "track_state.png",
        theme,
        caption="Rules are defined by `trackManagement.maintenance` and `trackManagement.deletion` in the config.",
    )

    # Slide 13: Track initiation + classification
    tm = cfg.get("trackManagement", {})
    init_cfg = (tm.get("initiation") or {})
    maint_cfg = (tm.get("maintenance") or {})
    del_cfg = (tm.get("deletion") or {})
    s = prs.slides.add_slide(blank)
    add_bullets(
        s,
        "Track initiation, maintenance, deletion, classification",
        [
            f"Initiation: M-of-N candidates (m={init_cfg.get('m','?')}, n={init_cfg.get('n','?')}) on unmatched clusters; optional initial velocity from two-point differencing.",
            f"Maintenance: confirmHits={maint_cfg.get('confirmHits','?')} and quality update (boost={maint_cfg.get('qualityBoost','?')}, decay={maint_cfg.get('qualityDecayRate','?')}).",
            f"Deletion: maxCoastingDwells={del_cfg.get('maxCoastingDwells','?')}, minQuality={del_cfg.get('minQuality','?')}, maxRange={del_cfg.get('maxRange','?')}.",
            "Classification: simple heuristics using speed + IMM mode probabilities (clutter/drone/bird/unknown).",
        ],
        theme,
    )

    # Slide 14: Configuration layout
    sys_cfg = cfg.get("system", {})
    net_cfg = cfg.get("network", {})
    cl_cfg = cfg.get("clustering", {})
    s = prs.slides.add_slide(blank)
    add_bullets(
        s,
        "Configuration (tracker_config.json) — key knobs",
        [
            f"System: cyclePeriodMs={sys_cfg.get('cyclePeriodMs','?')}  maxDetectionsPerDwell={sys_cfg.get('maxDetectionsPerDwell','?')}  maxTracks={sys_cfg.get('maxTracks','?')}",
            f"Network: receiver={net_cfg.get('receiverIp','?')}:{net_cfg.get('receiverPort','?')}  sender={net_cfg.get('senderIp','?')}:{net_cfg.get('senderPort','?')}",
            f"Clustering: method={cl_cfg.get('method','?')}  •  DBSCAN eps(range/az/el)={cl_cfg.get('dbscan', {})}",
            f"Association: method={assoc.get('method','?')}  gatingThreshold={assoc.get('gatingThreshold','?')}",
            "Operational practice: keep config under version control; changes should be validated with simulators/log replay.",
        ],
        theme,
    )

    # Slide 15: Reliability & performance
    s = prs.slides.add_slide(blank)
    add_bullets(
        s,
        "Reliability, performance, and real-time behavior",
        [
            "Concurrency: receiver thread deserializes UDP messages and pushes them into a guarded queue; processing thread consumes messages.",
            "Back-pressure risk: if input rate exceeds processing, the queue grows; consider bounded queues + drop/merge policy for production.",
            "Timing: receiver uses socket receive timeout; processing loop uses `cyclePeriodMs` wait-for pacing.",
            "Numerics: IMM update uses simplified covariance update for performance (Joseph form is noted as alternative).",
            "Logging: binary logs can be enabled for offline analysis; ensure disk I/O budget is acceptable in real-time runs.",
        ],
        theme,
    )

    # Slide 16: Software design principles / patterns
    s = prs.slides.add_slide(blank)
    add_bullets(
        s,
        "Software design principles & patterns applied",
        [
            "Separation of concerns: UDP transport (`UdpSocket`/`MessageSerializer`) is isolated from tracking algorithms.",
            "Single Responsibility: each class has a focused scope (receive, preprocess, cluster, associate, predict, manage tracks, send).",
            "Strategy pattern: `IClusterer` and `IAssociator` select behavior via configuration without changing call sites.",
            "Composition over inheritance: engines own strategy instances; pipeline composes subsystems.",
            "Configurability: algorithm choice and thresholds are externalized to `tracker_config.json`.",
        ],
        theme,
    )

    # Slide 17: Extensibility
    s = prs.slides.add_slide(blank)
    add_bullets(
        s,
        "Extensibility guide (how to add new algorithms)",
        [
            "New clusterer: implement `IClusterer::cluster()`, add to `ClusterEngine` switch, extend config parsing if needed.",
            "New associator: implement `IAssociator::associate()`, add to `AssociationEngine` switch, add config knobs.",
            "New motion model: implement `IMotionModel`, register in `IMMFilter` models_ array, extend transition matrix sizing if expanded.",
            "Add telemetry: extend `BinaryLogger` record types and emit at stage boundaries (raw/preproc/cluster/assoc/update).",
        ],
        theme,
    )

    # Slide 18: Build & run (developer workflow)
    s = prs.slides.add_slide(blank)
    add_bullets(
        s,
        "Build / run / verify",
        [
            "Build (Linux): mkdir -p build && cd build && cmake .. && cmake --build . -j",
            "Run tracker: ./build/cuas_tracker [config/tracker_config.json]",
            "Simulators: dsp_injector (inject/replay detections), display_module (receive/display tracks), log_extractor (analyze logs).",
            "Integration test idea: replay a known detection log and assert stable track count / state transitions / message throughput.",
        ],
        theme,
    )

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))


def main() -> int:
    render_diagrams()
    build_deck()
    print(f"Wrote: {OUT_PPTX}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

